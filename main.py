import os
import io
import json
import asyncio
from typing import Optional

import pdfplumber
from pptx import Presentation
from fastapi import FastAPI, UploadFile, File, WebSocket, WebSocketDisconnect, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
import google.generativeai as genai
from dotenv import load_dotenv

load_dotenv()

app = FastAPI(title="STRIVE API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
model = genai.GenerativeModel("gemini-1.5-flash")

sessions = {}


def parse_pdf(file_bytes: bytes) -> str:
    text = ""
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for i, page in enumerate(pdf.pages):
            page_text = page.extract_text()
            if page_text:
                text += "\n\n--- Slide " + str(i+1) + " ---\n"
                text += page_text
    return text.strip()


def parse_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    text = ""
    for i, slide in enumerate(prs.slides):
        text += "\n\n--- Slide " + str(i+1) + " ---\n"
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text += shape.text.strip() + "\n"
    return text.strip()


@app.get("/")
def root():
    return {"status": "STRIVE API running"}


@app.post("/upload")
async def upload_slides(file: UploadFile = File(...)):
    file_bytes = await file.read()
    filename = file.filename.lower()

    if filename.endswith(".pdf"):
        slide_text = parse_pdf(file_bytes)
    elif filename.endswith(".pptx") or filename.endswith(".ppt"):
        slide_text = parse_pptx(file_bytes)
    else:
        return JSONResponse(status_code=400, content={"error": "Only PDF and PPTX files are supported."})

    import hashlib, time
    session_id = hashlib.md5((filename + str(time.time())).encode()).hexdigest()[:12]

    topic_prompt = (
        "You are an academic AI tutor. A student has uploaded lecture slides.\n"
        "Return a JSON object with: title, slide_count, topics (list of 3-6), summary.\n\n"
        "Slide content:\n" + slide_text[:4000] + "\n\n"
        "Respond ONLY with valid JSON. No markdown."
    )

    topic_response = model.generate_content(topic_prompt)
    raw = topic_response.text.strip().replace("```json", "").replace("```", "").strip()

    try:
        meta = json.loads(raw)
    except:
        meta = {
            "title": file.filename,
            "slide_count": slide_text.count("--- Slide"),
            "topics": ["Topic 1", "Topic 2", "Topic 3"],
            "summary": "Lecture slides uploaded successfully."
        }

    sessions[session_id] = {
        "filename": file.filename,
        "slide_text": slide_text,
        "meta": meta
    }

    return {"session_id": session_id, "filename": file.filename, "meta": meta}


@app.post("/chat")
async def chat(session_id: str = Form(...), message: str = Form(...)):
    session = sessions.get(session_id)
    slide_context = session["slide_text"][:6000] if session else ""

    if slide_context:
        prompt = (
            "You are STRIVE, an expert AI academic tutor.\n"
            "LECTURE SLIDES:\n" + slide_context + "\n\n"
            "Student: " + message + "\n"
            "Respond as a helpful clear tutor."
        )
    else:
        prompt = (
            "You are STRIVE, an expert AI academic tutor.\n"
            "Student: " + message + "\n"
            "Respond helpfully."
        )

    response = model.generate_content(prompt)
    return {"reply": response.text}


@app.post("/lesson")
async def generate_lesson(session_id: str = Form(...), topic: Optional[str] = Form(None)):
    session = sessions.get(session_id)
    if not session:
        return JSONResponse(status_code=404, content={"error": "Session not found."})

    slide_text = session["slide_text"]
    topic_str = "Focus on: " + topic if topic else "Cover all topics."

    prompt = (
        "You are STRIVE, an expert AI academic tutor.\n"
        + topic_str + "\n"
        "Generate a clear lesson with headings, bold key terms, short paragraphs, and a summary.\n\n"
        "LECTURE SLIDES:\n" + slide_text[:8000]
    )

    response = model.generate_content(prompt)
    return {"lesson": response.text}


@app.post("/exercises")
async def generate_exercises(session_id: str = Form(...)):
    session = sessions.get(session_id)
    if not session:
        return JSONResponse(status_code=404, content={"error": "Session not found."})

    slide_text = session["slide_text"]

    prompt = (
        "Generate 5 practice exercises from these lecture slides.\n"
        "Return a JSON array. Each item: id, type (mcq or short_answer), "
        "question, options (4 items or null), correct, explanation.\n"
        "At least 3 MCQ and 2 short answer.\n\n"
        "SLIDES:\n" + slide_text[:6000] + "\n\n"
        "Respond ONLY with a JSON array. No markdown."
    )

    response = model.generate_content(prompt)
    raw = response.text.strip().replace("```json", "").replace("```", "").strip()

    try:
        exercises = json.loads(raw)
    except:
        exercises = []

    return {"exercises": exercises}


@app.post("/exam-questions")
async def generate_exam_questions(session_id: str = Form(...)):
    session = sessions.get(session_id)
    if not session:
        return JSONResponse(status_code=404, content={"error": "Session not found."})

    slide_text = session["slide_text"]

    prompt = (
        "Generate 5 likely university exam questions from these slides.\n"
        "Return a JSON array. Each item: id, type (essay/short_answer/mcq), "
        "question, marks, difficulty, marking_guide, options (or null).\n"
        "Mix: 2 essay, 2 short_answer, 1 mcq.\n\n"
        "SLIDES:\n" + slide_text[:6000] + "\n\n"
        "Respond ONLY with a JSON array. No markdown."
    )

    response = model.generate_content(prompt)
    raw = response.text.strip().replace("```json", "").replace("```", "").strip()

    try:
        questions = json.loads(raw)
    except:
        questions = []

    return {"exam_questions": questions}


@app.websocket("/ws/{session_id}")
async def websocket_chat(websocket: WebSocket, session_id: str):
    await websocket.accept()

    try:
        while True:
            message = await websocket.receive_text()
            session = sessions.get(session_id)
            slide_context = session["slide_text"][:6000] if session else ""

            if slide_context:
                prompt = (
                    "You are STRIVE, an expert AI academic tutor.\n"
                    "SLIDES:\n" + slide_context + "\n\nStudent: " + message
                )
            else:
                prompt = "You are STRIVE, an AI tutor.\nStudent: " + message

            response = model.generate_content(prompt, stream=True)

            for chunk in response:
                if chunk.text:
                    await websocket.send_json({"type": "token", "content": chunk.text})
                    await asyncio.sleep(0.01)

            await websocket.send_json({"type": "done"})

    except WebSocketDisconnect:
        print("Client disconnected: " + session_id)
    except Exception as e:
        await websocket.send_json({"type": "error", "content": str(e)})


if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)
